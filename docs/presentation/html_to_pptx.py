#!/usr/bin/env python3
"""
Script untuk screenshot slides HTML dan mengubahnya menjadi PPTX
Requires: pip install playwright python-pptx
Setup: playwright install chromium
"""

import os
import asyncio
from pathlib import Path

# Check and install dependencies
try:
    from playwright.async_api import async_playwright
except ImportError:
    print("Installing playwright...")
    os.system("pip install playwright")
    os.system("playwright install chromium")
    from playwright.async_api import async_playwright

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Installing python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt


async def screenshot_slides(html_path: str, output_dir: str) -> list:
    """Screenshot each slide from the HTML presentation"""
    
    screenshots = []
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        # Use 2x scale factor for sharper images (2560x1440 effective resolution)
        page = await browser.new_page(
            viewport={'width': 1280, 'height': 720},
            device_scale_factor=2  # 2x resolution for crisp screenshots
        )
        
        # Open the HTML file
        file_url = f"file://{os.path.abspath(html_path)}"
        await page.goto(file_url)
        await page.wait_for_load_state('networkidle')
        
        # Wait for fonts to load
        await page.wait_for_timeout(1000)
        
        # Find all slides
        slides = await page.query_selector_all('.slide')
        print(f"Found {len(slides)} slides")
        
        for i, slide in enumerate(slides, 1):
            screenshot_path = os.path.join(output_dir, f"slide_{i:02d}.png")
            
            # Scroll slide into view and screenshot
            await slide.scroll_into_view_if_needed()
            await asyncio.sleep(0.3)  # Wait for any animations
            await slide.screenshot(path=screenshot_path)
            
            screenshots.append(screenshot_path)
            print(f"Screenshot saved: {screenshot_path}")
        
        await browser.close()
    
    return screenshots


def create_pptx(screenshots: list, output_path: str):
    """Create a PPTX from screenshot images"""
    
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank slide layout
    blank_layout = prs.slide_layouts[6]
    
    for screenshot in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add image to fill the entire slide
        slide.shapes.add_picture(
            screenshot,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    prs.save(output_path)
    print(f"\nPPTX saved: {output_path}")


async def main():
    # Configuration
    script_dir = os.path.dirname(os.path.abspath(__file__))
    html_path = os.path.join(script_dir, "index.html")
    screenshots_dir = os.path.join(script_dir, "screenshots")
    pptx_path = os.path.join(script_dir, "Presentasi_Sidang_Skripsi.pptx")
    
    print("=" * 50)
    print("HTML to PPTX Converter")
    print("=" * 50)
    print(f"Input:  {html_path}")
    print(f"Output: {pptx_path}")
    print("=" * 50)
    
    # Step 1: Screenshot slides
    print("\n[1/2] Taking screenshots...")
    screenshots = await screenshot_slides(html_path, screenshots_dir)
    
    # Step 2: Create PPTX
    print("\n[2/2] Creating PPTX...")
    create_pptx(screenshots, pptx_path)
    
    print("\n✅ Done!")
    print(f"   PPTX file: {pptx_path}")
    print(f"   Screenshots: {screenshots_dir}/")


if __name__ == "__main__":
    asyncio.run(main())
